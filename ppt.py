#!/usr/bin/env python3
"""
build_ppt_from_results.py

Gera um arquivo PPTX (outputs/presentation.pptx) a partir dos resultados do script
`run_nlp_practice_no_ppt.py` — usa os arquivos em `outputs/figs/` e `outputs/experiment_results.csv`.

Instalação necessária:
    pip install python-pptx pandas

Como usar:
    python build_ppt_from_results.py

O script monta 5 slides (em português) conforme solicitado no practice:
 1. Resumo do experimento
 2. Curvas de aprendizado
 3. Comparação entre experimentos (tabela resumida)
 4. Visualização dos embeddings
 5. Conclusões e próximos passos

Coloca images automaticamente quando encontram arquivos na pasta de figuras.
"""

import os
import sys
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm

# Config
OUTDIR = Path("outputs")
FIGS_DIR = OUTDIR / "figs"
CSV_PATH = OUTDIR / "experiment_results.csv"
OUT_PPTX = OUTDIR / "presentation.pptx"

# Texts em português (já prontos)
TITLE = "Experimentos: Embeddings e Classificação de Texto"
SLIDE1_BULLETS = [
    "Objetivo: avaliar impacto de tamanho do vocabulário, tamanho do dataset e dimensão dos embeddings.",
    "Dataset: IMDb (subamostras conforme experimentos).",
    "Modelo: Embedding -> mean pooling -> classificação linear (PyTorch).",
    "Métricas: acurácia, loss, F1; visualizações: PCA / t-SNE dos embeddings."
]

SLIDE2_TITLE = "Curvas de aprendizado"
SLIDE2_TEXT = (
    "Mostramos as curvas de loss e accuracy (treino vs validação).\n"
    "Verifique convergência e sinais de overfitting/underfitting para cada configuração."
)

SLIDE3_TITLE = "Comparação quantitativa entre experimentos"
SLIDE3_TEXT = (
    "Resumo dos resultados agregados por configuração (accuracy e val_f1).\n"
    "A tabela abaixo mostra as métricas finais usadas para comparação."
)

SLIDE4_TITLE = "Visualização dos Embeddings"
SLIDE4_TEXT = (
    "Projetei os embeddings (PCA + t-SNE) e colorizei por classe.\n"
    "A separação espacial indica o quanto as classes estão discrimináveis."
)

SLIDE5_TITLE = "Conclusões e próximos passos"
SLIDE5_BULLETS = [
    "Aumentar vocab_size melhora até um ponto; custo computacional aumenta.",
    "Embedding dim maiores tendem a melhorar separabilidade, mas com retornos decrescentes.",
    "Dataset size tem forte impacto na generalização; mais dados ajudam mais que aumentar dim.",
    "Próximos passos: testar sequence models (GRU/LSTM) e regularização/augmentation."
]

# Helpers de layout
def add_title(slide, text):
    title = slide.shapes.title
    title.text = text

def add_bullets(slide, bullets, left=Inches(0.6), top=Inches(1.6), width=Inches(8), height=Inches(3.5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

def add_paragraph_text(slide, text, left=Inches(0.6), top=Inches(1.6), width=Inches(8), height=Inches(3.5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)

def insert_image_safe(slide, img_path, left, top, width=None, height=None):
    try:
        if width is not None and height is not None:
            slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)
        elif width is not None:
            slide.shapes.add_picture(str(img_path), left, top, width=width)
        elif height is not None:
            slide.shapes.add_picture(str(img_path), left, top, height=height)
        else:
            slide.shapes.add_picture(str(img_path), left, top)
        return True
    except Exception as e:
        print(f"Aviso: não foi possível inserir imagem {img_path}: {e}")
        return False


def find_figures_for_slide(keyword):
    """Procura figuras na pasta FIGS_DIR contendo a palavra-chave no nome (case-insensitive)."""
    if not FIGS_DIR.exists():
        return []
    k = keyword.lower()
    imgs = sorted([p for p in FIGS_DIR.glob("*.png") if k in p.name.lower()])
    return imgs


def build_presentation():
    prs = Presentation()

    # Slide 1: Resumo
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, TITLE)
    add_bullets(slide, SLIDE1_BULLETS)

    # Slide 2: Curvas de aprendizado
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, SLIDE2_TITLE)
    add_paragraph_text(slide, SLIDE2_TEXT, top=Inches(1.4))
    # Insert up to 2 representative curve images side by side
    curves = find_figures_for_slide("curve_")
    if len(curves) == 0:
        curves = find_figures_for_slide("curve")
    # choose up to 2
    for i, img in enumerate(curves[:2]):
        left = Inches(0.5 + i*5)
        insert_image_safe(slide, img, left, Inches(2.2), width=Inches(4.6))

    # Slide 3: Comparação (tabela simplificada)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, SLIDE3_TITLE)
    add_paragraph_text(slide, SLIDE3_TEXT, top=Inches(1.4))

    # Load CSV and build a small summary table (top 6 rows sorted by val_f1)
    if CSV_PATH.exists():
        try:
            df = pd.read_csv(CSV_PATH)
            if not df.empty:
                cols = ['experiment','vocab_size','embed_dim','dataset_size','val_acc','val_f1']
                present_cols = [c for c in cols if c in df.columns]
                df_tbl = df[present_cols].sort_values(by='val_f1' if 'val_f1' in df.columns else present_cols[-1], ascending=False).head(6)
                # Add table
                rows = df_tbl.shape[0] + 1
                cols_n = df_tbl.shape[1]
                left = Inches(0.5)
                top = Inches(2.2)
                width = Inches(9)
                height = Inches(1.6)
                table = slide.shapes.add_table(rows, cols_n, left, top, width, height).table
                # header
                for j, c in enumerate(df_tbl.columns):
                    table.cell(0, j).text = str(c)
                # rows
                for i, (_, r) in enumerate(df_tbl.iterrows(), start=1):
                    for j, c in enumerate(df_tbl.columns):
                        table.cell(i, j).text = f"{r[c]:.4f}" if isinstance(r[c], float) else str(r[c])
        except Exception as e:
            print("Aviso: falha ao carregar CSV para tabela:", e)
    else:
        add_paragraph_text(slide, "Arquivo de resultados não encontrado para gerar a tabela.", top=Inches(2.2))

    # Slide 4: Embeddings
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, SLIDE4_TITLE)
    add_paragraph_text(slide, SLIDE4_TEXT, top=Inches(1.4))
    emb_imgs = find_figures_for_slide("emb_")
    if len(emb_imgs) == 0:
        emb_imgs = find_figures_for_slide("emb")
    # insert up to 2 embedding plots
    for i, img in enumerate(emb_imgs[:2]):
        left = Inches(0.5 + i*5)
        insert_image_safe(slide, img, left, Inches(2.2), width=Inches(4.6))

    # Slide 5: Conclusões
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "Conclusões & Próximos passos")
    add_bullets(slide, SLIDE5_BULLETS)

    # Save
    ensure_outdir()
    prs.save(str(OUT_PPTX))
    print(f"Apresentação salva em: {OUT_PPTX}")


def ensure_outdir():
    OUTDIR.mkdir(parents=True, exist_ok=True)
    FIGS_DIR.mkdir(parents=True, exist_ok=True)


if __name__ == '__main__':
    build_presentation()
    print("Concluído.")
